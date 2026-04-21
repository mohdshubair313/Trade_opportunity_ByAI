"""
Report export service — turns a stored Analysis into PDF / XLSX / PPTX bytes.

Everything is pure-Python so it runs on python:3.11-slim without extra apt
packages. We trade a bit of visual polish for deployment simplicity; we can
swap in weasyprint later if the consultant persona demands magazine-quality
PDFs.
"""
from __future__ import annotations

import io
import re
from dataclasses import dataclass
from datetime import datetime
from typing import List, Optional, Sequence

from reportlab.lib import colors
from reportlab.lib.pagesizes import A4
from reportlab.lib.styles import getSampleStyleSheet, ParagraphStyle
from reportlab.lib.units import cm
from reportlab.platypus import (
    SimpleDocTemplate,
    Paragraph,
    Spacer,
    PageBreak,
    ListFlowable,
    ListItem,
)

from openpyxl import Workbook
from openpyxl.styles import Font, Alignment, PatternFill
from openpyxl.utils import get_column_letter

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor


# ---------------------------------------------------------------------------
# Shared markdown parser
# ---------------------------------------------------------------------------

@dataclass
class Section:
    """A single top-level section of the report (H2)."""
    title: str
    paragraphs: List[str]
    bullets: List[str]


_METADATA_FENCE = re.compile(r"^---\s*\n(.*?)\n---\s*\n", re.DOTALL)
_H1 = re.compile(r"^#\s+(.+)$")
_H2 = re.compile(r"^##\s+(.+)$")
_H3 = re.compile(r"^###\s+(.+)$")
_BULLET = re.compile(r"^\s*[-*]\s+(.+)$")
# Strip inline citation chips like [3] or [1][7] so exported copies stay readable.
_CITATION = re.compile(r"\s?\[\d+\]")


def _strip_frontmatter(report: str) -> str:
    match = _METADATA_FENCE.match(report)
    if not match:
        return report
    return report[match.end():]


def _clean(text: str) -> str:
    text = _CITATION.sub("", text)
    # Drop markdown bold/italic markers for plain-text exports.
    text = re.sub(r"\*\*(.+?)\*\*", r"\1", text)
    text = re.sub(r"\*(.+?)\*", r"\1", text)
    return text.strip()


def _parse_sections(report: str) -> List[Section]:
    body = _strip_frontmatter(report)
    sections: List[Section] = []
    current: Optional[Section] = None
    current_paragraph: List[str] = []

    def flush_paragraph():
        nonlocal current_paragraph
        if current and current_paragraph:
            current.paragraphs.append(_clean(" ".join(current_paragraph)))
        current_paragraph = []

    for raw in body.splitlines():
        line = raw.rstrip()
        if not line.strip():
            flush_paragraph()
            continue

        h1 = _H1.match(line)
        if h1:
            flush_paragraph()
            continue  # H1 is used as the document title; handled separately

        h2 = _H2.match(line)
        if h2:
            flush_paragraph()
            current = Section(title=_clean(h2.group(1)), paragraphs=[], bullets=[])
            sections.append(current)
            continue

        h3 = _H3.match(line)
        if h3:
            flush_paragraph()
            if current is None:
                current = Section(title="Overview", paragraphs=[], bullets=[])
                sections.append(current)
            current.paragraphs.append(_clean(h3.group(1)).upper())
            continue

        bullet = _BULLET.match(line)
        if bullet:
            flush_paragraph()
            if current is None:
                current = Section(title="Overview", paragraphs=[], bullets=[])
                sections.append(current)
            current.bullets.append(_clean(bullet.group(1)))
            continue

        if current is None:
            current = Section(title="Summary", paragraphs=[], bullets=[])
            sections.append(current)
        current_paragraph.append(_clean(line))

    flush_paragraph()
    return sections


def _extract_title(report: str) -> str:
    body = _strip_frontmatter(report)
    for line in body.splitlines():
        m = _H1.match(line.rstrip())
        if m:
            return _clean(m.group(1))
    return "Sector Analysis"


# ---------------------------------------------------------------------------
# PDF (reportlab)
# ---------------------------------------------------------------------------

def _pdf_styles():
    base = getSampleStyleSheet()
    return {
        "title": ParagraphStyle(
            "Title",
            parent=base["Title"],
            fontSize=20,
            leading=24,
            textColor=colors.HexColor("#0a7a3b"),
            spaceAfter=12,
        ),
        "meta": ParagraphStyle(
            "Meta",
            parent=base["Normal"],
            fontSize=9,
            textColor=colors.HexColor("#666666"),
            spaceAfter=18,
        ),
        "h2": ParagraphStyle(
            "H2",
            parent=base["Heading2"],
            fontSize=14,
            leading=18,
            textColor=colors.HexColor("#0a7a3b"),
            spaceBefore=16,
            spaceAfter=8,
        ),
        "body": ParagraphStyle(
            "Body",
            parent=base["BodyText"],
            fontSize=10.5,
            leading=15,
            textColor=colors.HexColor("#222222"),
            spaceAfter=8,
        ),
        "bullet": ParagraphStyle(
            "Bullet",
            parent=base["BodyText"],
            fontSize=10.5,
            leading=15,
            textColor=colors.HexColor("#222222"),
            leftIndent=14,
            bulletIndent=4,
        ),
        "source": ParagraphStyle(
            "Source",
            parent=base["BodyText"],
            fontSize=9,
            leading=12,
            textColor=colors.HexColor("#333333"),
            leftIndent=14,
            spaceAfter=4,
        ),
    }


def report_to_pdf(
    *,
    report: str,
    sector: str,
    sources_analyzed: int,
    generated_at: datetime,
    sources: Sequence[dict] = (),
) -> bytes:
    buffer = io.BytesIO()
    doc = SimpleDocTemplate(
        buffer,
        pagesize=A4,
        leftMargin=2 * cm,
        rightMargin=2 * cm,
        topMargin=2 * cm,
        bottomMargin=2 * cm,
        title=f"{sector.title()} Analysis",
        author="TradeInsight AI",
    )
    styles = _pdf_styles()
    story = []

    title = _extract_title(report) or f"{sector.title()} Sector Analysis"
    story.append(Paragraph(title, styles["title"]))
    story.append(
        Paragraph(
            f"Generated {generated_at.strftime('%d %B %Y, %H:%M UTC')} &middot; "
            f"{sources_analyzed} source(s) analyzed &middot; TradeInsight AI",
            styles["meta"],
        )
    )

    for section in _parse_sections(report):
        story.append(Paragraph(section.title, styles["h2"]))
        for para in section.paragraphs:
            if para:
                story.append(Paragraph(para, styles["body"]))
        if section.bullets:
            story.append(
                ListFlowable(
                    [ListItem(Paragraph(b, styles["bullet"])) for b in section.bullets],
                    bulletType="bullet",
                    start="•",
                )
            )
            story.append(Spacer(1, 6))

    if sources:
        story.append(PageBreak())
        story.append(Paragraph("Sources", styles["h2"]))
        for src in sources:
            n = src.get("n", "?")
            title_text = (src.get("title") or src.get("url") or "source").replace("&", "&amp;")
            url = src.get("url", "")
            line = f"<b>[{n}]</b> {title_text}<br/><font color='#0a7a3b'>{url}</font>"
            story.append(Paragraph(line, styles["source"]))

    doc.build(story)
    return buffer.getvalue()


# ---------------------------------------------------------------------------
# XLSX (openpyxl)
# ---------------------------------------------------------------------------

def report_to_xlsx(
    *,
    report: str,
    sector: str,
    sources_analyzed: int,
    generated_at: datetime,
    sources: Sequence[dict] = (),
) -> bytes:
    wb = Workbook()
    header_font = Font(bold=True, color="FFFFFF")
    header_fill = PatternFill(fill_type="solid", start_color="0A7A3B", end_color="0A7A3B")
    wrap = Alignment(wrap_text=True, vertical="top")

    # Overview sheet
    overview = wb.active
    overview.title = "Overview"
    overview["A1"] = f"{sector.title()} Sector Analysis"
    overview["A1"].font = Font(bold=True, size=14, color="0A7A3B")
    overview["A3"] = "Generated"
    overview["B3"] = generated_at.strftime("%d %B %Y %H:%M UTC")
    overview["A4"] = "Sector"
    overview["B4"] = sector
    overview["A5"] = "Sources analyzed"
    overview["B5"] = sources_analyzed
    for row in range(3, 6):
        overview.cell(row=row, column=1).font = Font(bold=True)
    overview.column_dimensions["A"].width = 22
    overview.column_dimensions["B"].width = 60

    # One section per row on the Summary sheet
    summary = wb.create_sheet("Summary")
    summary.append(["Section", "Key points"])
    for cell in summary[1]:
        cell.font = header_font
        cell.fill = header_fill
    for section in _parse_sections(report):
        points: List[str] = []
        points.extend(section.paragraphs)
        points.extend(f"• {b}" for b in section.bullets)
        summary.append([section.title, "\n".join(p for p in points if p)])
        summary.cell(row=summary.max_row, column=2).alignment = wrap
    summary.column_dimensions["A"].width = 32
    summary.column_dimensions["B"].width = 90
    for row in range(2, summary.max_row + 1):
        summary.row_dimensions[row].height = 80

    # Sources sheet
    if sources:
        src_sheet = wb.create_sheet("Sources")
        src_sheet.append(["#", "Title", "URL", "Snippet"])
        for cell in src_sheet[1]:
            cell.font = header_font
            cell.fill = header_fill
        for src in sources:
            src_sheet.append([
                src.get("n"),
                src.get("title") or "",
                src.get("url") or "",
                (src.get("snippet") or "")[:280],
            ])
        for idx, width in enumerate([6, 50, 60, 60], start=1):
            src_sheet.column_dimensions[get_column_letter(idx)].width = width

    # Raw markdown sheet (full fidelity)
    raw = wb.create_sheet("Raw Markdown")
    raw["A1"] = "Raw report (markdown)"
    raw["A1"].font = Font(bold=True)
    raw["A3"] = report
    raw["A3"].alignment = wrap
    raw.column_dimensions["A"].width = 120

    out = io.BytesIO()
    wb.save(out)
    return out.getvalue()


# ---------------------------------------------------------------------------
# PPTX (python-pptx)
# ---------------------------------------------------------------------------

_BRAND_GREEN = RGBColor(0x0A, 0x7A, 0x3B)
_MUTED = RGBColor(0x55, 0x55, 0x55)
_DARK = RGBColor(0x22, 0x22, 0x22)


def report_to_pptx(
    *,
    report: str,
    sector: str,
    sources_analyzed: int,
    generated_at: datetime,
    sources: Sequence[dict] = (),
) -> bytes:
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    blank = prs.slide_layouts[6]

    # Title slide
    slide = prs.slides.add_slide(blank)
    _add_textbox(slide, 0.75, 2.0, 11.8, 1.6,
                 f"{sector.title()} Sector Analysis",
                 size=44, bold=True, color=_BRAND_GREEN)
    _add_textbox(slide, 0.75, 3.6, 11.8, 0.8,
                 "AI-powered market intelligence",
                 size=20, color=_DARK)
    _add_textbox(slide, 0.75, 4.4, 11.8, 0.6,
                 f"Generated {generated_at.strftime('%d %B %Y')}  •  "
                 f"{sources_analyzed} sources  •  TradeInsight AI",
                 size=14, color=_MUTED)

    # One slide per H2 section (capped at 12 so decks stay skimmable)
    for section in _parse_sections(report)[:12]:
        slide = prs.slides.add_slide(blank)
        _add_textbox(slide, 0.6, 0.4, 12.1, 0.9,
                     section.title, size=28, bold=True, color=_BRAND_GREEN)

        # Body: first paragraph as lead + bullets
        body_top = 1.5
        tb = slide.shapes.add_textbox(Inches(0.6), Inches(body_top), Inches(12.1), Inches(5.4)).text_frame
        tb.word_wrap = True

        first = True
        for para in section.paragraphs:
            if not para:
                continue
            p = tb.paragraphs[0] if first else tb.add_paragraph()
            first = False
            run = p.add_run()
            run.text = para
            run.font.size = Pt(14)
            run.font.color.rgb = _DARK

        for bullet in section.bullets:
            p = tb.paragraphs[0] if first else tb.add_paragraph()
            first = False
            run = p.add_run()
            run.text = f"•  {bullet}"
            run.font.size = Pt(14)
            run.font.color.rgb = _DARK
            p.level = 0

    # Sources slide
    if sources:
        slide = prs.slides.add_slide(blank)
        _add_textbox(slide, 0.6, 0.4, 12.1, 0.9, "Sources",
                     size=28, bold=True, color=_BRAND_GREEN)
        tb = slide.shapes.add_textbox(Inches(0.6), Inches(1.5), Inches(12.1), Inches(5.4)).text_frame
        tb.word_wrap = True
        first = True
        for src in sources[:15]:
            p = tb.paragraphs[0] if first else tb.add_paragraph()
            first = False
            n = src.get("n", "?")
            title = src.get("title") or src.get("url") or "source"
            url = src.get("url") or ""
            run = p.add_run()
            run.text = f"[{n}] {title}"
            run.font.size = Pt(12)
            run.font.bold = True
            run.font.color.rgb = _DARK
            if url:
                url_para = tb.add_paragraph()
                url_run = url_para.add_run()
                url_run.text = url
                url_run.font.size = Pt(10)
                url_run.font.color.rgb = _MUTED

    out = io.BytesIO()
    prs.save(out)
    return out.getvalue()


def _add_textbox(slide, x: float, y: float, w: float, h: float,
                 text: str, *, size: int = 14, bold: bool = False,
                 color: Optional[RGBColor] = None) -> None:
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    if color is not None:
        run.font.color.rgb = color


# ---------------------------------------------------------------------------
# Public helpers
# ---------------------------------------------------------------------------

CONTENT_TYPES = {
    "pdf": "application/pdf",
    "xlsx": "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet",
    "pptx": "application/vnd.openxmlformats-officedocument.presentationml.presentation",
    "md": "text/markdown; charset=utf-8",
}


def export_analysis(
    *,
    fmt: str,
    report: str,
    sector: str,
    sources_analyzed: int,
    generated_at: datetime,
    sources: Sequence[dict] = (),
) -> bytes:
    fmt = fmt.lower()
    kwargs = dict(
        report=report,
        sector=sector,
        sources_analyzed=sources_analyzed,
        generated_at=generated_at,
        sources=sources,
    )
    if fmt == "pdf":
        return report_to_pdf(**kwargs)
    if fmt == "xlsx":
        return report_to_xlsx(**kwargs)
    if fmt == "pptx":
        return report_to_pptx(**kwargs)
    if fmt == "md":
        return report.encode("utf-8")
    raise ValueError(f"Unsupported format: {fmt}")
